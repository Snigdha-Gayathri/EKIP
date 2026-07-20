"""
Document Ingestion Agent — File Parsers and Chunkers
"""

from __future__ import annotations

import csv
import io
import logging
from typing import Any

from bs4 import BeautifulSoup
import docx
import pypdf

logger = logging.getLogger(__name__)


def parse_document(file_content: bytes, file_type: str, filename: str) -> str:
    """Parse raw bytes into plain text based on file type."""
    file_type = file_type.lower().lstrip(".")
    if file_type == "pdf":
        return _parse_pdf(file_content)
    elif file_type == "docx":
        return _parse_docx(file_content)
    elif file_type == "pptx":
        return _parse_pptx(file_content)
    elif file_type == "html":
        return _parse_html(file_content)
    elif file_type == "csv":
        return _parse_csv(file_content)
    else:
        # Default: treat as markdown or txt
        return file_content.decode("utf-8", errors="ignore")


def _parse_pdf(content: bytes) -> str:
    try:
        reader = pypdf.PdfReader(io.BytesIO(content))
        pages = []
        for idx, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                pages.append(f"--- Page {idx+1} ---\n{text}")
        if pages:
            return "\n\n".join(pages)
    except Exception as e:
        logger.warning("pypdf extraction failed (%s), trying unstructured fallback", e)
    try:
        from unstructured.partition.pdf import partition_pdf
        elements = partition_pdf(file=io.BytesIO(content))
        return "\n".join(str(el) for el in elements)
    except Exception as e:
        logger.warning("PDF unstructured fallback failed: %s", e)
        return ""


def _parse_docx(content: bytes) -> str:
    try:
        doc = docx.Document(io.BytesIO(content))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
    except Exception as e:
        logger.warning("docx extraction failed (%s), trying unstructured fallback", e)
    try:
        from unstructured.partition.docx import partition_docx
        elements = partition_docx(file=io.BytesIO(content))
        return "\n".join(str(el) for el in elements)
    except Exception as e:
        logger.warning("DOCX unstructured fallback failed: %s", e)
        return ""


def _parse_pptx(content: bytes) -> str:
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(content))
        text_runs = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                text_runs.append(f"--- Slide {idx+1} ---\n" + "\n".join(slide_text))
        if text_runs:
            return "\n\n".join(text_runs)
    except Exception as e:
        logger.warning("python-pptx extraction failed (%s), trying unstructured fallback", e)
    try:
        from unstructured.partition.pptx import partition_pptx
        elements = partition_pptx(file=io.BytesIO(content))
        return "\n".join(str(el) for el in elements)
    except Exception as e:
        logger.warning("PPTX unstructured fallback failed: %s", e)
        return ""


def _parse_html(content: bytes) -> str:
    soup = BeautifulSoup(content, "html.parser")
    for script in soup(["script", "style"]):
        script.extract()
    return soup.get_text(separator="\n")


def _parse_csv(content: bytes) -> str:
    text = content.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = [", ".join(row) for row in reader]
    return "\n".join(rows)


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[dict[str, Any]]:
    """Split text recursively into chunks with section title tracking."""
    lines = text.splitlines()
    chunks = []
    current_chunk: list[str] = []
    current_length = 0
    current_section = "Overview"

    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            current_section = stripped.lstrip("#").strip()

        line_len = len(line) + 1
        if current_length + line_len > chunk_size and current_chunk:
            chunk_str = "\n".join(current_chunk)
            chunks.append(
                {
                    "text": chunk_str,
                    "section_title": current_section,
                }
            )
            # Retain overlap lines
            overlap_chars = 0
            overlap_lines = []
            for l in reversed(current_chunk):
                if overlap_chars + len(l) <= overlap:
                    overlap_lines.insert(0, l)
                    overlap_chars += len(l)
                else:
                    break
            current_chunk = overlap_lines
            current_length = overlap_chars

        current_chunk.append(line)
        current_length += line_len

    if current_chunk:
        chunks.append(
            {
                "text": "\n".join(current_chunk),
                "section_title": current_section,
            }
        )

    return chunks
