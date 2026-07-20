import os
import sys
from pathlib import Path

def create_pdf(path: Path):
    # Create a minimal valid PDF containing text
    pdf_content = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 125 >>\nstream\n"
        b"BT\n/F1 12 Tf\n50 700 Td\n(Acme AI Enterprise Demo PDF Document) Tj\n0 -20 Td\n(Zero-trust authentication and microservice architecture guide.) Tj\nET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n0000000241 00000 n \n0000000418 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\nstartxref\n488\n%%EOF\n"
    )
    path.write_bytes(pdf_content)
    print(f"Created {path} ({len(pdf_content)} bytes)")

def create_docx(path: Path):
    try:
        import docx
        doc = docx.Document()
        doc.add_heading("Acme AI Enterprise Demo DOCX Document", 0)
        doc.add_paragraph("This document summarizes key microservices: api-gateway, auth-service, and billing-service.")
        doc.add_paragraph("All deployments on Render Free tier must utilize 1 uvicorn worker and lazy loading for embeddings.")
        doc.save(path)
        print(f"Created {path} using python-docx")
    except ImportError:
        import zipfile, io
        # Minimal valid docx zip structure
        out = io.BytesIO()
        with zipfile.ZipFile(out, "w") as z:
            z.writestr("[Content_Types].xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>')
            z.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>')
            z.writestr("word/document.xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Acme AI Enterprise Demo DOCX Document. Microservices architecture and zero-trust guidelines.</w:t></w:r></w:p></w:body></w:document>')
        path.write_bytes(out.getvalue())
        print(f"Created {path} using zip structure ({len(out.getvalue())} bytes)")

def create_pptx(path: Path):
    try:
        import pptx
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Acme AI Enterprise Demo PPTX"
        slide.placeholders[1].text = "Overview of microservices:\n- API Gateway (Rate limiting & mTLS)\n- Auth Service (JWT zero-trust)\n- Billing Service (Stripe integration)"
        prs.save(path)
        print(f"Created {path} using python-pptx")
    except ImportError:
        import zipfile, io
        out = io.BytesIO()
        with zipfile.ZipFile(out, "w") as z:
            z.writestr("[Content_Types].xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/><Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/></Types>')
            z.writestr("_rels/.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/></Relationships>')
            z.writestr("ppt/_rels/presentation.xml.rels", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/></Relationships>')
            z.writestr("ppt/presentation.xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:sldIdLst><p:sldId id="256" r:id="rId1" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"/></p:sldIdLst></p:presentation>')
            z.writestr("ppt/slides/slide1.xml", '<?xml version="1.0" encoding="UTF-8" standalone="yes"?><p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:sp><p:txBody><a:p><a:r><a:t>Acme AI Enterprise Demo PPTX. Slide 1 overview of microservices.</a:t></a:r></a:p></p:txBody></p:sp></p:spTree></p:cSld></p:sld>')
        path.write_bytes(out.getvalue())
        print(f"Created {path} using zip structure ({len(out.getvalue())} bytes)")

if __name__ == "__main__":
    demo_dir = Path(__file__).resolve().parent
    create_pdf(demo_dir / "sample.pdf")
    create_docx(demo_dir / "sample.docx")
    create_pptx(demo_dir / "sample.pptx")
